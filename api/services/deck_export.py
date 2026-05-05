"""
PPTX + DOCX export for presentations.

Direct python-pptx / python-docx generation, scoped to AHNJ's section types
(narrative / table / attachment / react_component). For react_component
sections we emit a placeholder note since we can't render TSX server-side
without a headless browser.

Adapted from the bank-processor's document_builder approach but slimmed
down to the section types AHNJ actually uses.
"""
from __future__ import annotations

import io
import logging
import os
import re
import tempfile
import uuid
from typing import Any, Optional

logger = logging.getLogger(__name__)

# Brand color matches the rest of the app (`#385854`).
AH_BRAND_RGB = (56, 88, 84)


# ─── Markdown → text segments (lightweight, pptx/docx-friendly) ───────────

_HEADING_RE = re.compile(r"^(#+)\s+(.+)$")
_BULLET_RE = re.compile(r"^\s*[-*]\s+(.+)$")
_NUM_RE = re.compile(r"^\s*\d+\.\s+(.+)$")
_CITATION_RE = re.compile(r"\[source:\s*([^\]]+)\]", re.IGNORECASE)


def _strip_citations(text: str) -> str:
    """Replace `[source: foo.pdf]` with `(foo.pdf)` so it shows in the export
    without the markdown shorthand looking out of place."""
    return _CITATION_RE.sub(r"(\1)", text or "")


def _markdown_blocks(md: str):
    """Yield (kind, text, level) tuples. kinds: heading, bullet, number, para.
    Tables in markdown are not parsed here — section-level table type already
    has structured headers/rows. This handles narrative bodies."""
    if not md:
        return
    paragraph_buf: list[str] = []

    def flush_para():
        nonlocal paragraph_buf
        if paragraph_buf:
            yield ("para", " ".join(paragraph_buf).strip(), 0)
            paragraph_buf = []

    for raw_line in md.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            yield from flush_para()
            continue
        m = _HEADING_RE.match(line)
        if m:
            yield from flush_para()
            yield ("heading", _strip_citations(m.group(2).strip()), len(m.group(1)))
            continue
        m = _BULLET_RE.match(line)
        if m:
            yield from flush_para()
            yield ("bullet", _strip_citations(m.group(1).strip()), 0)
            continue
        m = _NUM_RE.match(line)
        if m:
            yield from flush_para()
            yield ("number", _strip_citations(m.group(1).strip()), 0)
            continue
        paragraph_buf.append(_strip_citations(line.strip()))
    yield from flush_para()


# ─── DOCX export ──────────────────────────────────────────────────────────

def export_docx(presentation) -> tuple[bytes, str]:
    """Build a DOCX from the presentation, return (bytes, filename)."""
    from docx import Document as DocxDocument
    from docx.shared import Pt, RGBColor, Inches

    doc = DocxDocument()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    # Title page
    title = doc.add_heading(presentation.title or "Untitled Presentation", level=0)
    for run in title.runs:
        run.font.color.rgb = RGBColor(*AH_BRAND_RGB)
    if presentation.published_at:
        p = doc.add_paragraph(f"Published {presentation.published_at.strftime('%B %d, %Y')}")
        for run in p.runs:
            run.italic = True
            run.font.size = Pt(10)
    doc.add_paragraph()

    for section in (presentation.sections or []):
        kind = section.get("kind", "narrative")
        title_text = section.get("title")
        if title_text:
            h = doc.add_heading(title_text, level=1)
            for run in h.runs:
                run.font.color.rgb = RGBColor(*AH_BRAND_RGB)

        if kind == "narrative":
            for blk_kind, text, level in _markdown_blocks(section.get("body") or ""):
                if blk_kind == "heading":
                    doc.add_heading(text, level=min(max(level + 1, 2), 4))
                elif blk_kind == "bullet":
                    doc.add_paragraph(text, style="List Bullet")
                elif blk_kind == "number":
                    doc.add_paragraph(text, style="List Number")
                else:
                    doc.add_paragraph(text)

        elif kind == "table":
            headers = section.get("headers") or []
            rows = section.get("rows") or []
            if headers or rows:
                cols = max(len(headers), max((len(r) for r in rows), default=0)) or 1
                tbl = doc.add_table(rows=1 + len(rows), cols=cols)
                tbl.style = "Light Grid Accent 1"
                if headers:
                    hdr = tbl.rows[0].cells
                    for i in range(cols):
                        hdr[i].text = headers[i] if i < len(headers) else ""
                        for run in hdr[i].paragraphs[0].runs:
                            run.bold = True
                for ri, row in enumerate(rows):
                    cells = tbl.rows[ri + 1].cells
                    for ci in range(cols):
                        cells[ci].text = row[ci] if ci < len(row) else ""
            if section.get("caption"):
                cap = doc.add_paragraph(section["caption"])
                for run in cap.runs:
                    run.italic = True
                    run.font.size = Pt(9)

        elif kind == "attachment":
            att_id = section.get("attachment_id")
            atts = presentation.attachments or []
            att = next((a for a in atts if a.get("id") == att_id), None)
            label = (att or {}).get("filename") or att_id or "(attachment)"
            p = doc.add_paragraph()
            p.add_run("[Attachment: ").italic = True
            run = p.add_run(label)
            run.bold = True
            p.add_run("]").italic = True
            if section.get("caption"):
                cap = doc.add_paragraph(section["caption"])
                for run in cap.runs:
                    run.italic = True
                    run.font.size = Pt(9)

        elif kind == "react_component":
            p = doc.add_paragraph()
            p.add_run("[Custom component — render in the web viewer for full visualization.]").italic = True

        doc.add_paragraph()  # blank line between sections

    buf = io.BytesIO()
    doc.save(buf)
    safe_title = re.sub(r"[^\w\-]+", "-", (presentation.title or "presentation")).strip("-")[:80]
    return buf.getvalue(), f"{safe_title or 'presentation'}.docx"


# ─── PPTX export ──────────────────────────────────────────────────────────

def export_pptx(presentation) -> tuple[bytes, str]:
    """Build a PPTX from the presentation, return (bytes, filename)."""
    from pptx import Presentation as Pptx

    deck = Pptx()
    blank_layout = deck.slide_layouts[6]   # entirely blank
    title_layout = deck.slide_layouts[0]   # title slide

    # Title slide
    s = deck.slides.add_slide(title_layout)
    s.shapes.title.text = presentation.title or "Untitled Presentation"
    if presentation.published_at:
        s.placeholders[1].text = presentation.published_at.strftime("%B %d, %Y")

    for section in (presentation.sections or []):
        kind = section.get("kind", "narrative")

        if kind == "narrative":
            _add_narrative_slides(deck, section, blank_layout)
        elif kind == "table":
            _add_table_slide(deck, section, blank_layout)
        elif kind == "attachment":
            _add_attachment_slide(deck, section, presentation.attachments or [], blank_layout)
        elif kind == "react_component":
            _add_react_placeholder_slide(deck, section, blank_layout)

    buf = io.BytesIO()
    deck.save(buf)
    safe_title = re.sub(r"[^\w\-]+", "-", (presentation.title or "presentation")).strip("-")[:80]
    return buf.getvalue(), f"{safe_title or 'presentation'}.pptx"


def _add_narrative_slides(deck, section: dict, layout):
    """Split the narrative across multiple slides if it's long. Each slide
    holds the section title + a chunk of the body."""
    from pptx.util import Inches, Pt

    title = section.get("title") or ""
    blocks = list(_markdown_blocks(section.get("body") or ""))

    # Naive chunking: 6 blocks per slide
    PER_SLIDE = 6
    chunks = [blocks[i:i + PER_SLIDE] for i in range(0, len(blocks), PER_SLIDE)] or [[]]

    for i, chunk in enumerate(chunks):
        slide = deck.slides.add_slide(layout)
        # Title text box
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title + (f" ({i + 1}/{len(chunks)})" if len(chunks) > 1 else "")
            run.font.size = Pt(28)
            run.font.bold = True
            from pptx.dml.color import RGBColor
            run.font.color.rgb = RGBColor(*AH_BRAND_RGB)

        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        first = True
        for blk_kind, text, _level in chunk:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            if blk_kind == "heading":
                run = p.add_run()
                run.text = text
                run.font.size = Pt(20)
                run.font.bold = True
            elif blk_kind in ("bullet", "number"):
                p.level = 0
                marker = "• " if blk_kind == "bullet" else ""
                run = p.add_run()
                run.text = f"{marker}{text}"
                run.font.size = Pt(14)
            else:
                run = p.add_run()
                run.text = text
                run.font.size = Pt(14)


def _add_table_slide(deck, section: dict, layout):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = deck.slides.add_slide(layout)
    title = section.get("title")
    if title:
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(28); run.font.bold = True
        run.font.color.rgb = RGBColor(*AH_BRAND_RGB)

    headers = section.get("headers") or []
    rows = section.get("rows") or []
    cols = max(len(headers), max((len(r) for r in rows), default=0))
    if cols == 0:
        return

    tbl_rows = 1 + len(rows)
    table_shape = slide.shapes.add_table(tbl_rows, cols, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5))
    table = table_shape.table

    # header row
    for ci in range(cols):
        cell = table.cell(0, ci)
        cell.text = headers[ci] if ci < len(headers) else ""
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(11)

    # data rows
    for ri, row in enumerate(rows):
        for ci in range(cols):
            cell = table.cell(ri + 1, ci)
            cell.text = row[ci] if ci < len(row) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)

    if section.get("caption"):
        cb = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
        cp = cb.text_frame.paragraphs[0]
        run = cp.add_run()
        run.text = section["caption"]
        run.font.italic = True
        run.font.size = Pt(10)


def _add_attachment_slide(deck, section: dict, attachments: list, layout):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = deck.slides.add_slide(layout)
    att_id = section.get("attachment_id")
    att = next((a for a in (attachments or []) if a.get("id") == att_id), None)
    label = (att or {}).get("filename") or att_id or "(attachment)"

    # Title
    if section.get("title"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = section["title"]
        run.font.size = Pt(28); run.font.bold = True
        run.font.color.rgb = RGBColor(*AH_BRAND_RGB)

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4))
    tf = body.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Attachment"
    run.font.size = Pt(20); run.font.italic = True

    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = label
    r.font.size = Pt(16); r.font.bold = True

    if section.get("caption"):
        p3 = tf.add_paragraph()
        r = p3.add_run()
        r.text = section["caption"]
        r.font.size = Pt(12); r.font.italic = True


def _add_react_placeholder_slide(deck, section: dict, layout):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide = deck.slides.add_slide(layout)
    if section.get("title"):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = section["title"]
        run.font.size = Pt(28); run.font.bold = True
        run.font.color.rgb = RGBColor(*AH_BRAND_RGB)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(2))
    p = body.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "[Custom component — view in the web viewer for full visualization.]"
    run.font.italic = True
    run.font.size = Pt(14)
